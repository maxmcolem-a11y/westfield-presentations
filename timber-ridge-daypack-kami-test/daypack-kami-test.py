#!/usr/bin/env python3
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent
DAYPACK = ROOT.parent / 'westfield-presentations' / 'timber-ridge-daypack-development-clean' / 'assets'
OUT = ROOT / 'assets' / 'demos' / 'timber-ridge-daypack-kami-test.pptx'

PARCHMENT = RGBColor(0xF5, 0xF4, 0xED)
IVORY = RGBColor(0xFA, 0xF9, 0xF5)
BRAND = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK = RGBColor(0x14, 0x14, 0x13)
DARK_WARM = RGBColor(0x3D, 0x3D, 0x3A)
OLIVE = RGBColor(0x5E, 0x5D, 0x59)
STONE = RGBColor(0x87, 0x86, 0x7F)
BORDER = RGBColor(0xE8, 0xE6, 0xDC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SERIF = 'Newsreader'
SANS = 'Inter'
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def blank_slide(prs, bg_color=PARCHMENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_paragraphs(slide, items, left, top, width, height, bullet_indent=18):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = SANS
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_WARM
        p.level = 0
        p.bullet = True
        p.line_spacing = 1.25
        p.space_after = Pt(10)
        p.left_margin = Pt(bullet_indent)
        p.hanging = Pt(8)
    return tb


def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_card(slide, left, top, width, height, fill=IVORY, border=BORDER, border_weight=0.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card


def add_label(slide, text, left, top, width=Inches(2.2), fill=IVORY):
    tag = add_card(slide, left, top, width, Inches(0.45), fill=fill, border=BORDER, border_weight=0.75)
    add_text(slide, text, left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.22),
             font=SANS, size=10, bold=True, color=BRAND)
    return tag


def add_picture(slide, path, left, top, width=None, height=None):
    kwargs = {}
    if width is not None:
        kwargs['width'] = width
    if height is not None:
        kwargs['height'] = height
    return slide.shapes.add_picture(str(path), left, top, **kwargs)


def cover_slide(prs):
    s = blank_slide(prs)
    logo = DAYPACK / 'timber-ridge-logo-cover.png'
    if logo.exists():
        add_picture(s, logo, Inches(11.35), Inches(0.65), width=Inches(1.25))
    add_text(s, 'Timber Ridge 26L Daypack', Inches(0.95), Inches(2.0), Inches(11.4), Inches(0.9),
             font=SERIF, size=30, color=BRAND, align=PP_ALIGN.CENTER)
    add_text(s, 'Kami format test deck', Inches(1.0), Inches(2.85), Inches(11.3), Inches(0.8),
             font=SERIF, size=46, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.15), Inches(4.0), Inches(1.05), weight_pt=1.25)
    add_text(s,
             'Recasting the approved daypack development story into the warm parchment editorial slide system from tw93/kami.',
             Inches(1.3), Inches(4.35), Inches(10.7), Inches(0.9), font=SANS, size=17,
             color=OLIVE, align=PP_ALIGN.CENTER)
    add_text(s, 'Westfield Outdoors · Test build · 2026.04', Inches(1), Inches(6.7), Inches(11.3), Inches(0.3),
             font=SANS, size=12, color=STONE, align=PP_ALIGN.CENTER)


def toc_slide(prs):
    s = blank_slide(prs)
    add_text(s, 'Contents', Inches(1.0), Inches(0.7), Inches(3.5), Inches(0.6), font=SERIF, size=32)
    add_line(s, Inches(1.0), Inches(1.55), Inches(11.1), weight_pt=1)
    items = [
        'Mandate and market position',
        'Volumetric exploration in CAD',
        'Three finalists selected from 1/5 scale prints',
        'Diamond route, signature read and geometry',
        'Spiral route, wearability and access logic',
        'Takeaways from the Kami translation test',
    ]
    for i, item in enumerate(items):
        y = Inches(2.05 + i * 0.73)
        add_text(s, f'{i+1:02d}', Inches(1.05), y, Inches(0.7), Inches(0.35), font=SERIF, size=24, color=BRAND)
        add_text(s, item, Inches(2.0), y + Inches(0.02), Inches(9.8), Inches(0.45), font=SERIF, size=20, color=NEAR_BLACK)


def mandate_slide(prs):
    s = blank_slide(prs)
    add_label(s, '01 · Strategy', Inches(1.0), Inches(0.65))
    add_text(s, 'Mandate and position', Inches(1.0), Inches(1.3), Inches(5.8), Inches(0.8), font=SERIF, size=30)
    add_text(s,
             'The original deck opens with a tight product brief. In Kami form, that same idea reads more like an editorial thesis page.',
             Inches(1.0), Inches(2.1), Inches(5.3), Inches(0.9), font=SANS, size=16, color=OLIVE)
    add_paragraphs(s, [
        'Design an iconic, highly structured 26L silhouette anchored by a single signature element, recognizable from 20 feet away.',
        'Capture the $79 to $100 designed-middle segment with a premium visual read, not a commodity outdoor look.',
        'Build brand identity through function-driven geometry rather than decorative overlays or tactical cues.',
    ], Inches(1.05), Inches(3.0), Inches(5.1), Inches(2.7))
    add_card(s, Inches(7.0), Inches(1.35), Inches(5.1), Inches(4.8))
    img = DAYPACK / 'faceted-pillow-hero-1.jpg'
    if img.exists():
        add_picture(s, img, Inches(7.25), Inches(1.62), width=Inches(4.6), height=Inches(3.35))
    add_text(s, 'Editorial restraint fits the opening slide well. It makes the concept feel more premium, but it also reduces the punch of the original highly visual cover.',
             Inches(7.35), Inches(5.18), Inches(4.35), Inches(0.8), font=SANS, size=13, color=DARK_WARM)


def exploration_slide(prs):
    s = blank_slide(prs)
    add_label(s, '02 · Form Study', Inches(1.0), Inches(0.65))
    add_text(s, 'Initial volumetric exploration', Inches(1.0), Inches(1.25), Inches(6), Inches(0.8), font=SERIF, size=30)
    add_text(s,
             'The first phase used 3D modeling to generate and compare distinct volumetric directions for recognizability, tension, proportion, and structural clarity.',
             Inches(1.0), Inches(2.0), Inches(11.0), Inches(0.8), font=SANS, size=17, color=OLIVE)
    img = DAYPACK / 'TR_Daypack_CADs_W5.2.png'
    if img.exists():
        add_picture(s, img, Inches(1.0), Inches(2.85), width=Inches(7.8), height=Inches(3.85))
    add_card(s, Inches(9.15), Inches(2.85), Inches(3.1), Inches(3.85))
    add_text(s, 'What Kami changes', Inches(9.45), Inches(3.15), Inches(2.4), Inches(0.3), font=SANS, size=11, bold=True, color=STONE)
    add_paragraphs(s, [
        'Turns the exploration sheet into a calmer museum-wall presentation.',
        'Helps the CAD grid feel intentional and collectible.',
        'Feels best when paired with one sharp explanatory sentence, not a lot of bullets.',
    ], Inches(9.42), Inches(3.55), Inches(2.3), Inches(2.4), bullet_indent=14)


def finalists_slide(prs):
    s = blank_slide(prs)
    add_label(s, '03 · Selection', Inches(1.0), Inches(0.65))
    add_text(s, 'Three shapes advanced after physical review', Inches(1.0), Inches(1.2), Inches(8), Inches(0.8), font=SERIF, size=30)
    add_text(s, 'All modeled shapes were printed at 1/5 scale, then reviewed internally for tactile viability and immediate visual appeal.',
             Inches(1.0), Inches(1.95), Inches(10.8), Inches(0.6), font=SANS, size=16, color=OLIVE)
    cards = [
        ('Diamond', DAYPACK / 'shape-1-3d-print.jpg'),
        ('Teardrop', DAYPACK / 'shape-2-3d-print.jpg'),
        ('Spiral', DAYPACK / 'shape-3-3d-print.jpg'),
    ]
    x_positions = [Inches(0.95), Inches(4.45), Inches(7.95)]
    for x, (label, path) in zip(x_positions, cards):
        add_card(s, x, Inches(2.7), Inches(3.05), Inches(3.55))
        if path.exists():
            add_picture(s, path, x + Inches(0.16), Inches(2.9), width=Inches(2.73), height=Inches(2.55))
        add_text(s, label, x + Inches(0.18), Inches(5.62), Inches(1.7), Inches(0.3), font=SERIF, size=22, color=NEAR_BLACK)
        add_text(s, 'Advanced to concept refinement', x + Inches(0.18), Inches(5.98), Inches(2.5), Inches(0.3), font=SANS, size=12, color=OLIVE)


def diamond_slide(prs):
    s = blank_slide(prs)
    add_label(s, '04 · Diamond Route', Inches(1.0), Inches(0.65), width=Inches(2.5))
    add_text(s, 'Diamond final design', Inches(1.0), Inches(1.25), Inches(5.4), Inches(0.7), font=SERIF, size=30)
    add_text(s, 'A highly structured silhouette with strong graphic potential and function-led brand identity.',
             Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.7), font=SANS, size=16, color=OLIVE)
    add_paragraphs(s, [
        'Triangular panels establish the immediate structured read.',
        'The lower M-shape webbing projects an inner V upward to create the slanted zipper access point.',
        'Laptop storage and quick-access organization stay embedded within the form, rather than interrupting it.',
    ], Inches(1.03), Inches(2.9), Inches(5.0), Inches(2.8))
    hero = DAYPACK / 'faceted-pillow-hero-2.jpg'
    if hero.exists():
        add_picture(s, hero, Inches(7.0), Inches(1.5), width=Inches(5.15), height=Inches(4.15))
    add_text(s, 'Kami read: this route translates especially well because the bag already has a disciplined, graphic silhouette.',
             Inches(7.15), Inches(5.95), Inches(4.8), Inches(0.5), font=SANS, size=13, color=DARK_WARM)


def spiral_slide(prs):
    s = blank_slide(prs)
    add_label(s, '05 · Spiral Route', Inches(1.0), Inches(0.65), width=Inches(2.4))
    add_text(s, 'Spiral final design', Inches(1.0), Inches(1.25), Inches(5.4), Inches(0.7), font=SERIF, size=30)
    add_text(s, 'Hybrid sling and pack functionality driven by spiral panel architecture.',
             Inches(1.0), Inches(2.0), Inches(5.3), Inches(0.7), font=SANS, size=16, color=OLIVE)
    add_paragraphs(s, [
        'The panel flow bridges commuter access, ergonomic body wrap, and an athletic identity.',
        'Spiral panels create a recognizable motion language from multiple angles.',
        'The asymmetrical belt helps merge sling immediacy with backpack stability.',
    ], Inches(1.02), Inches(2.95), Inches(5.0), Inches(2.7))
    hero = DAYPACK / 'spiral-hero-1.jpg'
    if hero.exists():
        add_picture(s, hero, Inches(7.0), Inches(1.5), width=Inches(5.15), height=Inches(4.15))
    add_text(s, 'Kami read: the spiral route feels more experimental and fashion-forward, but it wants generous image space to really sell the motion.',
             Inches(7.15), Inches(5.92), Inches(4.75), Inches(0.6), font=SANS, size=13, color=DARK_WARM)


def comparison_slide(prs):
    s = blank_slide(prs)
    add_label(s, '06 · Test Takeaways', Inches(1.0), Inches(0.65), width=Inches(2.7))
    add_text(s, 'What this format does well for the daypack story', Inches(1.0), Inches(1.25), Inches(8.4), Inches(0.8), font=SERIF, size=30)
    cols = [
        ('Strengths', [
            'Elevates the development story into something that feels more editorial and investor-facing.',
            'Works especially well for strategy, thesis, and winner-selection slides.',
            'Makes the Diamond concept feel premium and disciplined.',
        ]),
        ('Tradeoffs', [
            'Less expressive than the current custom HTML deck for zoomable image storytelling.',
            'Not as good for dense process slides with many comparative images.',
            'Would likely want fewer slides, each with stronger hero imagery and shorter text.',
        ]),
    ]
    x_positions = [Inches(1.0), Inches(6.75)]
    for x, (title, bullets) in zip(x_positions, cols):
        add_card(s, x, Inches(2.3), Inches(5.05), Inches(3.7))
        add_text(s, title, x + Inches(0.25), Inches(2.6), Inches(2.8), Inches(0.4), font=SERIF, size=24, color=BRAND)
        add_paragraphs(s, bullets, x + Inches(0.24), Inches(3.15), Inches(4.5), Inches(2.4), bullet_indent=14)
    add_line(s, Inches(1.0), Inches(6.55), Inches(11.1), weight_pt=1)
    add_text(s, 'Recommendation: use Kami as a polished executive-summary companion deck, not a full replacement for the richer interactive HTML development deck.',
             Inches(1.0), Inches(6.72), Inches(11.2), Inches(0.4), font=SANS, size=13, color=STONE, align=PP_ALIGN.CENTER)


def ending_slide(prs):
    s = blank_slide(prs)
    add_text(s, 'Kami daypack test complete', Inches(1.0), Inches(3.0), Inches(11.33), Inches(0.8), font=SERIF, size=40, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.2), Inches(1.0), weight_pt=1.5)
    add_text(s, 'Next step: if you like the direction, I can build a cleaner neutral Westfield executive deck template around it.',
             Inches(1.2), Inches(4.55), Inches(10.9), Inches(0.7), font=SANS, size=16, color=OLIVE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    cover_slide(prs)
    toc_slide(prs)
    mandate_slide(prs)
    exploration_slide(prs)
    finalists_slide(prs)
    diamond_slide(prs)
    spiral_slide(prs)
    comparison_slide(prs)
    ending_slide(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f'✓ Saved {OUT}')


if __name__ == '__main__':
    main()
